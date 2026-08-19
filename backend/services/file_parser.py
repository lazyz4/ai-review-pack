"""从上传文件中提取纯文本。

支持 PPTX / PDF / DOCX / TXT / Markdown；旧版 .ppt 会提示另存为 .pptx。
提取结果用于后续 Ollama 分析与知识点生成。
"""

from __future__ import annotations

from io import BytesIO
from pathlib import Path

SUPPORTED_EXTENSIONS = {".pptx", ".pdf", ".docx", ".txt", ".md"}


def extract_text(filename: str, content: bytes) -> tuple[str, str]:
    """按扩展名提取文本，返回 (文本, 来源类型)。"""
    extension = Path(filename or "").suffix.lower()
    if extension == ".pptx":
        text = _extract_pptx(content)
    elif extension == ".pdf":
        text = _extract_pdf(content)
    elif extension == ".docx":
        text = _extract_docx(content)
    elif extension in {".txt", ".md"}:
        text = _extract_text_file(content)
    elif extension == ".ppt":
        raise ValueError("暂不支持旧版 .ppt 二进制格式，请在 PowerPoint 中另存为 .pptx 后重试")
    else:
        raise ValueError(f"不支持的文件类型：{extension or '未知'}。支持：PPTX / PDF / DOCX / TXT / Markdown")
    return text.strip(), extension.lstrip(".")


def _extract_pptx(content: bytes) -> str:
    """提取 PPTX 中文本框与表格的文字，按幻灯片分组。"""
    from pptx import Presentation

    presentation = Presentation(BytesIO(content))
    parts: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        lines = [f"--- 幻灯片 {index} ---"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in paragraph.runs).strip()
                    if text:
                        lines.append(text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    lines.append(" | ".join(cell.text.strip() for cell in row.cells))
        parts.append("\n".join(lines))
    return "\n\n".join(parts)


def _extract_pdf(content: bytes) -> str:
    """提取 PDF 每页文本。"""
    from pypdf import PdfReader

    reader = PdfReader(BytesIO(content))
    parts: list[str] = []
    for index, page in enumerate(reader.pages, start=1):
        text = (page.extract_text() or "").strip()
        parts.append(f"--- 第 {index} 页 ---\n{text}" if text else f"--- 第 {index} 页（无文本）---")
    return "\n\n".join(parts)


def _extract_docx(content: bytes) -> str:
    """提取 DOCX 段落与表格文字。"""
    from docx import Document

    document = Document(BytesIO(content))
    parts = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
    for table in document.tables:
        for row in table.rows:
            parts.append(" | ".join(cell.text.strip() for cell in row.cells))
    return "\n".join(parts)


def _extract_text_file(content: bytes) -> str:
    """按 UTF-8 / GBK / UTF-16 顺序尝试解码文本文件。"""
    for encoding in ("utf-8", "gbk", "utf-16"):
        try:
            return content.decode(encoding)
        except UnicodeDecodeError:
            continue
    return content.decode("utf-8", errors="replace")
